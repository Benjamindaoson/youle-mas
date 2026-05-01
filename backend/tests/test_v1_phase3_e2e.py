"""V1 Phase 3 端到端 — 声明式 skill + I/V/D agents。

DEMO 模式下能跑通的：
- monthly_report_pptx：T 写大纲 → D 出真实 .pptx 文件
- ecommerce_main_image：T 写图 prompt → I 走 placeholder（无 SiliconFlow key）→ T 写 alt
- run_skill 声明式路径：未带 runner 的 skill 也能执行

需要外部 API / FFmpeg 才能完整跑的（不在本测试覆盖）：
- V agent 完整音视频合成（依赖 ffmpeg + minimax）
- I agent 真实 AI 出图（依赖 SiliconFlow）
"""
from __future__ import annotations

import os

import pytest

from app.skills.registry import (
    SkillSpec,
    SkillStep,
    get_skill,
    load_all,
    run_skill,
)


@pytest.fixture(autouse=True, scope="module")
def _bootstrap_skills():
    n = load_all()
    assert n >= 3, "至少 3 个 skill：xiaohongshu / ecommerce_main_image / monthly_report_pptx + anti_scam_video"
    yield


# ===================== 声明式 skill 路径 =====================

async def _drain(spec: SkillSpec, message: str, session_id: str) -> list[dict]:
    out: list[dict] = []
    async for ev in run_skill(spec, {"message": message, "session_id": session_id}):
        out.append(ev)
    return out


async def test_run_skill_declarative_text_only_xiaohongshu():
    """xiaohongshu_hook_title 是单步 T，run_skill 声明式路径应能跑完。"""
    spec = get_skill("xiaohongshu_hook_title")
    assert spec is not None
    assert spec.runner is None  # 纯声明式

    events = await _drain(spec, "帮我给面膜写小红书标题", "v1:p3-xhs")
    types = [e["type"] for e in events]
    assert "skill_started" in types
    assert "skill_done" in types
    # 至少一个 T 产出
    assert any(e.get("type") == "artifact"
               and e.get("capability") == "T" for e in events)


async def test_run_skill_declarative_three_step_ecommerce():
    """ecommerce_main_image 三步：T → I → T。无 SiliconFlow 时 I 走 placeholder。"""
    spec = get_skill("ecommerce_main_image")
    assert spec is not None
    assert spec.runner is None

    events = await _drain(spec,
                          "帮我做一张面膜的电商主图,日系简约",
                          "v1:p3-ecom")
    types = [e["type"] for e in events]
    assert types.count("agent_start") == 3, \
        f"应有 3 步 agent_start，实际：{types}"
    assert types.count("agent_done") == 3
    # I agent 至少出一张图
    img_arts = [e for e in events
                if e.get("type") == "artifact" and e.get("capability") == "I"]
    assert img_arts, "I agent 至少要出一张图（placeholder fallback）"
    # 图片文件实际存在
    for art in img_arts:
        assert os.path.isfile(art["file_path"])


# ===================== D agent: PPT 真生成 =====================

async def test_d_agent_generates_real_pptx():
    """monthly_report_pptx：T 写大纲 → D 生成 .pptx 文件。"""
    spec = get_skill("monthly_report_pptx")
    assert spec is not None

    events = await _drain(spec,
                          "做一份咖啡店 8 月复盘",
                          "v1:p3-ppt")

    types = [e["type"] for e in events]
    assert "skill_done" in types
    assert "error" not in types, f"事件流不应有错误：{events}"

    # D agent 必须出一个 pptx artifact
    pptx_arts = [e for e in events
                 if e.get("type") == "artifact"
                 and e.get("artifact_type") == "pptx"]
    assert pptx_arts, "应当有 pptx artifact"
    pptx_path = pptx_arts[0]["file_path"]
    assert os.path.isfile(pptx_path), f"pptx 文件应存在：{pptx_path}"
    assert os.path.getsize(pptx_path) > 5000, "正常 pptx 至少几 KB"

    # 用 python-pptx 反读验证至少有 2 页
    from pptx import Presentation
    prs = Presentation(pptx_path)
    assert len(prs.slides) >= 2, f"PPT 至少要 2 页，实际 {len(prs.slides)}"


# ===================== V agent: 部分跑通（DEMO 不依赖 ffmpeg） =====================

async def test_v_agent_runs_with_fallbacks():
    """构造一个临时 V skill 验证 V agent 能 yield 出 voice + bgm + thumbnail/video-fallback。"""
    spec = SkillSpec(
        id="_test_video_only",
        name="V agent 单步测试",
        deliverable_type="video",
        steps=[SkillStep(agent="V",
                          task="为用户主题生成 30 秒短视频 (DEMO)")],
    )

    events: list[dict] = []
    async for ev in run_skill(spec, {
        "message": "做一个面膜介绍短视频",
        "session_id": "v1:p3-video",
    }):
        events.append(ev)

    arts = [e for e in events if e.get("type") == "artifact"]
    arts_kinds = {a.get("artifact_type") for a in arts}
    # 至少出 voice + bgm，video 或 video-fallback 二选一
    assert "voice" in arts_kinds, f"V agent 必出 voice，实际：{arts_kinds}"
    assert "bgm" in arts_kinds
    assert ("video" in arts_kinds) or ("video-fallback" in arts_kinds)


# ===================== Empty-steps skill 错误路径 =====================

async def test_run_skill_empty_steps_yields_error():
    spec = SkillSpec(id="_empty", name="空 skill", steps=[])
    events: list[dict] = []
    async for ev in run_skill(spec, {"message": "x", "session_id": "v1:p3-empty"}):
        events.append(ev)
    assert any(e.get("type") == "error" for e in events)
