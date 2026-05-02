"""D agent 可用的工具。

让"办公文档"能力 agent 不只是"出 PPT/Word"，还能"读 PDF / 透视 Excel /
看 PPT 大纲"。对应 agent4 = PDF / Excel / PPT。

V1 提供 4 个工具：
    - pdf_extract     : pypdf 抽前 N 页文本
    - excel_pivot     : pandas 看表格 schema + 简单透视（前 K 行 / dtypes / numeric 总和）
    - pptx_outline    : python-pptx 抽每页标题 + bullet
    - csv_pivot       : 同 excel_pivot 但只支持 csv（避免误用扩展名）

工具失败永不抛异常，返回 error dict。
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from app.config import settings
from app.logging_config import logger


TOOL_DEFS: list[dict[str, Any]] = [
    {
        "name": "pdf_extract",
        "description": (
            "用 pypdf 抽 PDF 前 N 页纯文本（不含图片/表格）。返回每页 ≤ 2000 字。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string"},
                "max_pages": {"type": "integer", "default": 10,
                              "description": "最多抽几页，1-50"},
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "excel_pivot",
        "description": (
            "读 .xlsx / .xls，返回列名 / dtypes / 行数 / 前 K 行 / 数值列总和。"
            "适用于先看表后再写分析报告。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string"},
                "sheet": {"type": "string",
                          "description": "工作表名；不填取第 1 个", "default": ""},
                "max_rows": {"type": "integer", "default": 20},
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "csv_pivot",
        "description": "同 excel_pivot 但读 .csv。",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string"},
                "max_rows": {"type": "integer", "default": 20},
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "pptx_outline",
        "description": "抽 .pptx 每页 title + 文字 bullets，返回 JSON list。",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string"},
                "max_slides": {"type": "integer", "default": 30},
            },
            "required": ["file_path"],
        },
    },
]


# ============================ 实现 ============================


def _validate_path(file_path: str, allowed_ext: tuple[str, ...]) -> tuple[str | None, dict | None]:
    abs_path = os.path.abspath(file_path)
    allowed_roots = [
        os.path.abspath(settings.UPLOAD_DIR),
        os.path.abspath(settings.ARTIFACT_DIR),
    ]
    if not any(abs_path.startswith(r) for r in allowed_roots):
        return None, {"error": "path_not_allowed", "file_path": file_path}
    if not os.path.isfile(abs_path):
        return None, {"error": "file_not_found", "file_path": file_path}
    ext = Path(abs_path).suffix.lower()
    if allowed_ext and ext not in allowed_ext:
        return None, {"error": "wrong_extension", "ext": ext,
                       "allowed": list(allowed_ext)}
    return abs_path, None


async def pdf_extract(file_path: str, max_pages: int = 10) -> dict:
    max_pages = max(1, min(int(max_pages or 10), 50))
    abs_path, err = _validate_path(file_path, (".pdf",))
    if err:
        return err
    try:
        from pypdf import PdfReader
    except ImportError:
        return {"error": "pypdf_missing"}

    try:
        reader = PdfReader(abs_path)
        n = len(reader.pages)
        pages = []
        for i, page in enumerate(reader.pages[:max_pages]):
            try:
                txt = (page.extract_text() or "").strip()
            except Exception as e:  # noqa: BLE001
                txt = f"[页 {i + 1} 抽取失败：{e}]"
            pages.append({"page": i + 1, "text": txt[:2000]})
    except Exception as e:  # noqa: BLE001
        return {"error": "pdf_parse_failed", "reason": str(e)}

    return {"file_path": file_path, "total_pages": n,
            "extracted_pages": len(pages), "pages": pages}


async def excel_pivot(file_path: str, sheet: str = "",
                       max_rows: int = 20) -> dict:
    max_rows = max(1, min(int(max_rows or 20), 200))
    abs_path, err = _validate_path(file_path, (".xlsx", ".xls"))
    if err:
        return err
    try:
        import pandas as pd
    except ImportError:
        return {"error": "pandas_missing"}

    try:
        if sheet:
            df = pd.read_excel(abs_path, sheet_name=sheet, nrows=max_rows + 50)
        else:
            df = pd.read_excel(abs_path, nrows=max_rows + 50)
    except Exception as e:  # noqa: BLE001
        return {"error": "excel_read_failed", "reason": str(e)}

    return _summarize_df(df, file_path, max_rows)


async def csv_pivot(file_path: str, max_rows: int = 20) -> dict:
    max_rows = max(1, min(int(max_rows or 20), 200))
    abs_path, err = _validate_path(file_path, (".csv",))
    if err:
        return err
    try:
        import pandas as pd
        df = pd.read_csv(abs_path, nrows=max_rows + 50)
    except ImportError:
        return {"error": "pandas_missing"}
    except Exception as e:  # noqa: BLE001
        return {"error": "csv_read_failed", "reason": str(e)}
    return _summarize_df(df, file_path, max_rows)


def _summarize_df(df, file_path: str, max_rows: int) -> dict:
    """共用的 df 摘要逻辑。"""
    head = df.head(max_rows)
    dtypes = {str(c): str(t) for c, t in df.dtypes.items()}
    numeric_sum: dict[str, float] = {}
    for col, t in df.dtypes.items():
        if str(t).startswith(("int", "float")):
            try:
                numeric_sum[str(col)] = float(df[col].sum())
            except (TypeError, ValueError):
                pass
    return {
        "file_path": file_path,
        "columns": [str(c) for c in head.columns],
        "row_count": int(df.shape[0]),
        "dtypes": dtypes,
        "numeric_sum": numeric_sum,
        "rows": head.to_dict(orient="records"),
    }


async def pptx_outline(file_path: str, max_slides: int = 30) -> dict:
    max_slides = max(1, min(int(max_slides or 30), 100))
    abs_path, err = _validate_path(file_path, (".pptx",))
    if err:
        return err
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python_pptx_missing"}

    try:
        prs = Presentation(abs_path)
    except Exception as e:  # noqa: BLE001
        return {"error": "pptx_open_failed", "reason": str(e)}

    slides_out: list[dict] = []
    for i, slide in enumerate(prs.slides):
        if i >= max_slides:
            break
        title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for j, para in enumerate(shape.text_frame.paragraphs):
                txt = "".join(r.text for r in para.runs).strip()
                if not txt:
                    continue
                if j == 0 and not title:
                    title = txt
                else:
                    bullets.append(txt)
        slides_out.append({"slide": i + 1, "title": title[:120],
                            "bullets": [b[:200] for b in bullets[:10]]})

    return {"file_path": file_path, "total_slides": len(prs.slides),
            "extracted_slides": len(slides_out), "slides": slides_out}


# ============================ Dispatch ============================


_TOOL_DISPATCH = {
    "pdf_extract": pdf_extract,
    "excel_pivot": excel_pivot,
    "csv_pivot": csv_pivot,
    "pptx_outline": pptx_outline,
}


async def call_tool(name: str, tool_input: dict) -> dict:
    fn = _TOOL_DISPATCH.get(name)
    if fn is None:
        return {"error": "unknown_tool", "name": name}
    try:
        return await fn(**(tool_input or {}))
    except TypeError as e:
        return {"error": "bad_arguments", "name": name, "reason": str(e)}
    except Exception as e:  # noqa: BLE001
        logger.warning("doc_tool_call_failed", name=name, error=str(e))
        return {"error": "tool_failed", "name": name, "reason": str(e)}


__all__ = ["TOOL_DEFS", "call_tool"]
