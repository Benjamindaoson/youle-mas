"""D agent — 办公文档能力（PDF / Excel / PPT / Word）。

V0 完全没有这个 capability。Phase 3 MVP：
- PPT 生成（python-pptx）
- DOCX 生成（python-docx）
- Excel 读（pandas，已有依赖）
- PDF 解析（pypdf；生成走 reportlab 留 V1.5）

任务路由：task.outputs 含 "ppt" / "pptx" → 出 PPT；"docx" → 出 Word；
其余默认出 markdown 报告（兜底）。
"""
from __future__ import annotations

import os
import re
import uuid
from collections.abc import AsyncIterator
from typing import Any, TYPE_CHECKING

from app.config import settings

from app.logging_config import logger
from app.skills.registry import SkillStep



if TYPE_CHECKING:
    from app.conductor.intent import Intent

def _safe(sid: str) -> str:
    s = re.sub(r"[^\w\-]", "_", sid or "")
    return re.sub(r"_+", "_", s).strip("_") or "default"


def _save_dir(session_id: str) -> str:
    base = os.path.join(settings.ARTIFACT_DIR, _safe(session_id), "docs")
    os.makedirs(base, exist_ok=True)
    return base


def _outputs_lower(task: SkillStep) -> str:
    return " ".join(task.outputs).lower() + " " + (task.task or "").lower()


def _slides_from_intent(intent: Intent, upstream: list[Any]) -> list[dict]:
    """构造 PPT 的 slides 数据。优先取上游 T agent 的 markdown 内容当大纲。"""
    title = intent.subject or "Youle 自动汇报"
    subtitle = intent.raw_user_text or ""

    body_text = ""
    for entry in upstream or []:
        for art in (entry.get("artifacts") or []):
            t = art.get("content_inline")
            if t:
                body_text += t + "\n\n"

    # 把 markdown 简单切成 slides：每个二级/无序列表当一页
    slides: list[dict] = [{"title": title, "subtitle": subtitle, "body": []}]
    cur: dict | None = None
    for line in (body_text or f"# {title}\n\n- 关于 {title} 的初步思考\n").splitlines():
        line = line.rstrip()
        if not line:
            continue
        if line.startswith("# ") or line.startswith("## "):
            cur = {"title": line.lstrip("# ").strip(), "subtitle": "", "body": []}
            slides.append(cur)
        elif line.startswith("- ") or line.startswith("* "):
            (cur or slides[-1])["body"].append(line[2:].strip())
        else:
            (cur or slides[-1])["body"].append(line.strip())
    if len(slides) == 1:
        # 没有结构 → 把 body_text 整段塞到一页
        slides.append({"title": title, "subtitle": "",
                       "body": [body_text[:300] or "（占位内容）"]})
    return slides


def _make_pptx(path: str, slides: list[dict]) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # 首页
    head = slides[0]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = head.get("title", "Untitled")
    if s0.placeholders and len(s0.placeholders) > 1:
        s0.placeholders[1].text = head.get("subtitle", "")

    for slide in slides[1:]:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide.get("title", "")
        body_holder = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                body_holder = ph
                break
        if body_holder is not None:
            tf = body_holder.text_frame
            tf.text = slide["body"][0] if slide.get("body") else ""
            for line in slide.get("body", [])[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

    prs.save(path)


def _make_docx(path: str, slides: list[dict]) -> None:
    from docx import Document
    doc = Document()
    head = slides[0]
    doc.add_heading(head.get("title", "Untitled"), level=0)
    if head.get("subtitle"):
        doc.add_paragraph(head["subtitle"])
    for slide in slides[1:]:
        doc.add_heading(slide.get("title", ""), level=1)
        for line in slide.get("body", []):
            doc.add_paragraph(line, style="List Bullet")
    doc.save(path)


_REACT_KEYWORDS = ("读 pdf", "读pdf", "解析", "抽取", "看表", "透视",
                    "读 excel", "读excel", "看 ppt", "看ppt", "outline",
                    "extract", "summarize", "分析表", "理解文档")


def _has_doc_in_upstream(upstream: list[Any]) -> bool:
    for entry in upstream or []:
        for art in (entry.get("artifacts") or []):
            atype = (art.get("artifact_type") or "").lower()
            if atype in ("pdf", "pptx", "docx", "xlsx", "csv"):
                return True
            fp = (art.get("file_path") or "").lower()
            if fp.endswith((".pdf", ".pptx", ".xlsx", ".xls", ".csv")):
                return True
    return False


def _should_use_react(task: SkillStep, upstream: list[Any]) -> bool:
    """task 含理解关键词 + 上游有 PDF/Excel/PPT + has_anthropic → ReAct。"""
    if not settings.has_anthropic:
        return False
    text = (task.task or "").lower()
    if not any(kw in text for kw in _REACT_KEYWORDS):
        return False
    return _has_doc_in_upstream(upstream)


async def run(
    task: SkillStep,
    intent: Intent,
    upstream: list[Any],
    session_id: str,
) -> AsyncIterator[dict]:
    """两条路径：
    - **ReAct 读文档**：上游有 PDF/Excel/PPT + task 含理解关键词
    - **生成**（默认）：根据 task.outputs 生成 .pptx / .docx / .md
    """
    if _should_use_react(task, upstream):
        try:
            async for ev in _react_loop(task=task, upstream=upstream,
                                         session_id=session_id):
                yield ev
            return
        except Exception as e:  # noqa: BLE001
            logger.warning("d_agent_react_failed_fallback_generate", error=str(e))

    save_dir = _save_dir(session_id)
    out_kind = _outputs_lower(task)
    slides = _slides_from_intent(intent, upstream)

    yield {"type": "chunk", "capability": "D",
           "text": f"准备生成文档（{len(slides)} 页）..."}

    try:
        if "pptx" in out_kind or "ppt" in out_kind:
            path = os.path.join(save_dir, f"deck_{uuid.uuid4().hex[:8]}.pptx")
            _make_pptx(path, slides)
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            artifact_type = "pptx"
        elif "docx" in out_kind or "word" in out_kind:
            path = os.path.join(save_dir, f"doc_{uuid.uuid4().hex[:8]}.docx")
            _make_docx(path, slides)
            mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            artifact_type = "docx"
        else:
            # 兜底：写 markdown
            path = os.path.join(save_dir, f"report_{uuid.uuid4().hex[:8]}.md")
            content = f"# {slides[0].get('title', '')}\n\n"
            for slide in slides[1:]:
                content += f"## {slide.get('title', '')}\n\n"
                for line in slide.get("body", []):
                    content += f"- {line}\n"
                content += "\n"
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            mime = "text/markdown"
            artifact_type = "markdown"
    except Exception as e:  # noqa: BLE001
        logger.error("d_agent_doc_failed", error=str(e), out_kind=out_kind)
        yield {"type": "error", "message": f"D agent 生成失败：{e}"}
        return

    yield {
        "type": "artifact",
        "capability": "D",
        "artifact_type": artifact_type,
        "title": task.task or "办公文档",
        "file_path": path,
        "mime_type": mime,
        "session_id": session_id,
    }
    yield {"type": "chunk", "capability": "D",
           "text": f"完成：{os.path.basename(path)}"}


# ============================ ReAct 路径（读 PDF/Excel/PPT）============================


_MAX_TOOL_TURNS = 4


async def _react_loop(*, task: SkillStep, upstream: list[Any],
                       session_id: str) -> AsyncIterator[dict]:
    """tool_use 循环：让 D agent 自主选 pdf_extract / excel_pivot / pptx_outline。"""
    import anthropic  # noqa: WPS433
    from app.adapters.model_router import pick_chat
    from app.capabilities.doc_tools import TOOL_DEFS, call_tool

    choice = pick_chat(purpose="capability_D_extract", prefer_provider="anthropic")
    if not choice.available:
        raise RuntimeError("anthropic not available for D ReAct")

    # 列出上游文档让 LLM 选
    doc_files: list[tuple[str, str]] = []  # (path, ext)
    for entry in upstream or []:
        for art in (entry.get("artifacts") or []):
            fp = art.get("file_path") or ""
            ext = os.path.splitext(fp)[1].lower()
            if ext in (".pdf", ".pptx", ".xlsx", ".xls", ".csv"):
                doc_files.append((fp, ext))

    file_listing = "\n".join(f"  - {p}  ({e})" for p, e in doc_files[:6])
    prompt = (
        f"任务：{task.task}\n\n"
        f"上游文档（{len(doc_files)} 个）：\n{file_listing}\n\n"
        f"你可以用 pdf_extract（读 PDF）、excel_pivot/csv_pivot（读表格）、"
        f"pptx_outline（读 PPT 大纲）。完成后给一段中文总结（≤400 字）。"
    )

    client = anthropic.AsyncAnthropic(api_key=choice.api_key)
    messages: list[dict[str, Any]] = [{"role": "user", "content": prompt}]

    yield {"type": "chunk", "capability": "D", "text": "[D ReAct] 启动文档分析..."}

    summary_text = ""
    for turn in range(_MAX_TOOL_TURNS + 1):
        tools_arg = TOOL_DEFS if turn < _MAX_TOOL_TURNS else None
        kwargs: dict[str, Any] = {
            "model": choice.model,
            "max_tokens": choice.max_tokens,
            "messages": messages,
        }
        if tools_arg:
            kwargs["tools"] = tools_arg

        resp = await client.messages.create(**kwargs)
        tool_uses: list[dict[str, Any]] = []
        round_text = ""
        for block in resp.content:
            if block.type == "text":
                round_text += block.text
            elif block.type == "tool_use":
                tool_uses.append({"id": block.id, "name": block.name,
                                   "input": dict(block.input or {})})

        if round_text:
            summary_text += round_text
            yield {"type": "chunk", "capability": "D", "text": round_text}

        if not tool_uses:
            break

        messages.append({"role": "assistant", "content": resp.content})

        tool_results: list[dict[str, Any]] = []
        for tu in tool_uses:
            yield {"type": "tool_call", "capability": "D",
                   "tool": tu["name"], "input": tu["input"], "turn": turn + 1}
            result = await call_tool(tu["name"], tu["input"])
            yield {"type": "tool_result", "capability": "D",
                   "tool": tu["name"], "result": result, "turn": turn + 1}
            tool_results.append({
                "type": "tool_result",
                "tool_use_id": tu["id"],
                "content": _serialize_tool_result(result),
            })
        messages.append({"role": "user", "content": tool_results})

    yield {
        "type": "artifact",
        "capability": "D",
        "artifact_type": "markdown",
        "title": task.task or "D 文档分析",
        "content_inline": summary_text or "（无文本输出）",
        "session_id": session_id,
    }


def _serialize_tool_result(result: Any) -> str:
    import json
    try:
        return json.dumps(result, ensure_ascii=False)[:8000]
    except (TypeError, ValueError):
        return str(result)[:8000]
